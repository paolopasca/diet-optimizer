#!/usr/bin/env python3
"""Genera il PPTX della dieta in stile 'diet template' (crema/verde) dall'output del solver.
Legge /tmp/week_plan.json (prodotto SEMPRE da make_week.py -> solve_diet.py).
Immagini in /tmp/dieta_img/. Output: ~/Desktop/dieta_paolo.pptx
"""
import json, os
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from PIL import Image

HERE = os.path.dirname(os.path.abspath(__file__))
IMGDIR = os.environ.get("DIETA_IMGDIR", "/tmp/dieta_img")  # esegui prima fetch_images.py
OUTDIR = os.environ.get("DIETA_OUTDIR", os.path.expanduser("~/Desktop/dieta"))
OUT = os.path.join(OUTDIR, "dieta.pptx")
os.makedirs(OUTDIR, exist_ok=True)
PLAN = json.load(open("/tmp/week_plan.json"))
DB = {f["name"]: f for f in json.load(open(os.path.join(HERE, "..", "data", "foods.json")))["foods"]}

# palette
CREAM = RGBColor(0xF3, 0xEC, 0xDD)
CARD = RGBColor(0xFF, 0xFD, 0xF8)
GREEN = RGBColor(0x5B, 0x6E, 0x35)
GREEN_D = RGBColor(0x3C, 0x49, 0x24)
ORANGE = RGBColor(0xD9, 0x84, 0x2B)
INK = RGBColor(0x32, 0x2E, 0x26)
GREY = RGBColor(0x8A, 0x82, 0x72)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
FONT = "Avenir Next"  # se assente, PowerPoint sostituisce con un sans pulito

EMU_IN = 914400
SW, SH = 13.333, 7.5

DAY_IMG = {"Lunedi": "fish", "Martedi": "chicken1", "Mercoledi": "burrito", "Giovedi": "greek",
           "Venerdi": "beef1", "Sabato": "caprese2", "Domenica": "legumi"}
DAY_TAG = {"Lunedi": "Pesce", "Martedi": "Pollo + tzatziki", "Mercoledi": "Burrito + tzatziki",
           "Giovedi": "Greca + gamberi", "Venerdi": "Tagliata di manzo",
           "Sabato": "Sgarro a cena", "Domenica": "Legumi + pesce"}
MEAL_ORDER = ["colazione", "pranzo", "merenda", "cena"]
MEAL_LABEL = {"colazione": "COLAZIONE", "pranzo": "PRANZO", "merenda": "MERENDA", "cena": "CENA"}

SHORT = {
 "Latte di vacca, UHT, parzialmente scremato": "latte p.s.",
 "Yogurt greco, 0% lipidi": "yogurt greco 0%", "Yogurt greco, da latte intero": "yogurt greco",
 "Corn flakes": "corn flakes", "Olio di oliva extra vergine": "olio EVO",
 "Pane di tipo integrale": "pane integrale", "Pasta di semola, integrale": "pasta integrale",
 "Pasta di semola": "pasta", "Tonno in salamoia, sgocciolato": "tonno al naturale",
 "Pollo, petto, crudo": "pollo", "Bovino adulto o vitellone, costata, crudo": "manzo (tagliata)",
 "Macinato di bovino magro (~5% grassi)": "macinato magro", "Orata, filetti": "orata",
 "Merluzzo o nasello": "merluzzo", "Gamberi sgusciati, surgelati": "gamberi",
 "Prosciutto crudo di Parma DOP, sgrassato": "prosciutto crudo",
 "Bresaola della Valtellina IGP": "bresaola", "Mozzarella di vacca": "mozzarella", "Feta": "feta",
 "Uova di gallina, intero": "uova", "Ceci, in scatola, scolati": "ceci",
 "Lenticchie, in scatola, scolate": "lenticchie", "Fagiolini, freschi": "fagiolini",
 "Pomodorini ciliegino, freschi": "pomodorini", "Pomodori, da insalata, freschi": "pomodori",
 "Cetrioli, freschi": "cetrioli", "Zucchine, chiare": "zucchine",
 "Melanzane, cotte, in padella": "melanzane", "Peperoni, crudi": "peperoni",
 "Lattuga, fresca": "lattuga", "Rucola o Ruchetta, fresca": "rucola",
 "Banane, fresche": "banana", "Pesche, fresche, con buccia": "pesca", "Cocomero, fresco": "cocomero",
 "Uva, fresca": "uva", "Mirtilli, freschi": "mirtilli", "Albicocche, fresche": "albicocche",
 "Avocado, fresco": "avocado", "Mandorle dolci, secche": "mandorle", "Noci, secche": "noci",
 "Miele": "miele", "Gelato confezionato, fior di latte, in vaschetta": "gelato fiordilatte",
}
# frutti grandi: plurale
PLUR = {"banana": "banane", "pesca": "pesche", "mela": "mele", "pera": "pere", "arancia": "arance"}

def short(n):
    return SHORT.get(n, n.split("(")[0].split(",")[0].strip().lower())

def label(it):
    """Etichetta completa: frutti grandi e uova a pezzi (col nome dentro), il resto nome+grammi."""
    nm = short(it["name"])
    sv = it.get("servings")
    if nm in PLUR and sv:
        return f"{sv} {nm if sv == 1 else PLUR[nm]}"
    if "uova" in it["name"].lower() and sv:
        return "1 uovo" if sv == 1 else f"{sv} uova"
    low = it["name"].lower()
    if "tonno" in low and sv:
        return f"{nm} {sv} {'scatoletta' if sv == 1 else 'scatolette'}"
    if ("ceci" in low or "lenticchie" in low) and "scatola" in low and sv:
        return f"{nm} {sv} {'barattolo' if sv == 1 else 'barattoli'}"
    return f"{nm} {int(round(it['grams']))}g"

def cover_crop(key, ratio, outname):
    src = os.path.join(IMGDIR, key + ".jpg")
    im = Image.open(src).convert("RGB")
    w, h = im.size
    cur = w / h
    if cur > ratio:  # troppo largo -> taglia ai lati
        nw = int(h * ratio); x = (w - nw) // 2; im = im.crop((x, 0, x + nw, h))
    else:            # troppo alto -> taglia sopra/sotto
        nh = int(w / ratio); y = (h - nh) // 2; im = im.crop((0, y, w, y + nh))
    out = os.path.join(IMGDIR, outname); im.save(out, quality=88); return out

prs = Presentation()
prs.slide_width = Inches(SW); prs.slide_height = Inches(SH)
BLANK = prs.slide_layouts[6]

def bg(slide, color=CREAM):
    s = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, Inches(SW), Inches(SH))
    s.fill.solid(); s.fill.fore_color.rgb = color; s.line.fill.background()
    s.shadow.inherit = False
    slide.shapes._spTree.remove(s._element); slide.shapes._spTree.insert(2, s._element)
    return s

def box(slide, l, t, w, h, color, line=None, rounded=False):
    shp = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE if rounded else MSO_SHAPE.RECTANGLE,
        Inches(l), Inches(t), Inches(w), Inches(h))
    shp.fill.solid(); shp.fill.fore_color.rgb = color
    if line: shp.line.color.rgb = line; shp.line.width = Pt(1)
    else: shp.line.fill.background()
    shp.shadow.inherit = False
    return shp

def text(slide, l, t, w, h, runs, align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP, sp=2):
    tb = slide.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h))
    tf = tb.text_frame; tf.word_wrap = True; tf.vertical_anchor = anchor
    tf.margin_left = tf.margin_right = Pt(4); tf.margin_top = tf.margin_bottom = Pt(2)
    if isinstance(runs, str): runs = [(runs, 18, False, INK)]
    first = True
    for item in runs:
        txt, size, bold, color = item
        p = tf.paragraphs[0] if first else tf.add_paragraph()
        first = False
        p.alignment = align; p.space_after = Pt(sp); p.space_before = Pt(0)
        r = p.add_run(); r.text = txt
        r.font.size = Pt(size); r.font.bold = bold; r.font.color.rgb = color; r.font.name = FONT
    return tb

def pic(slide, key, l, t, w, h, tag):
    p = cover_crop(key, w / h, f"_{tag}.jpg")
    slide.shapes.add_picture(p, Inches(l), Inches(t), Inches(w), Inches(h))

def badge(slide, l, t, w, label, value, color):
    box(slide, l, t, w, 0.95, color, rounded=True)
    text(slide, l, t + 0.12, w, 0.45, [(value, 21, True, WHITE)], PP_ALIGN.CENTER)
    text(slide, l, t + 0.58, w, 0.3, [(label, 9.5, True, WHITE)], PP_ALIGN.CENTER)

# ---------- COPERTINA ----------
s = prs.slides.add_slide(BLANK); bg(s)
pic(s, "hero", 7.0, 0, 6.333, 7.5, "cover")
box(s, 0, 0, 7.0, 7.5, CREAM)
text(s, 0.7, 1.5, 6.0, 1.0, [("La tua dieta", 46, True, GREEN_D)])
text(s, 0.7, 2.5, 6.0, 0.6, [("Estiva · Ricomposizione", 22, False, ORANGE)])
text(s, 0.72, 3.4, 5.8, 2.0, [
    ("~2400 kcal al giorno · 7 giorni", 16, True, INK),
    ("Proteine 150 g · Carboidrati ~300 g · Grassi ~67 g", 14, False, GREY),
    ("Sabato: pasto sgarro libero a cena", 14, False, GREY),
], sp=6)
box(s, 0.72, 5.6, 5.6, 1.1, CARD, rounded=True)
text(s, 0.95, 5.72, 5.2, 0.9, [
    ("Numeri reali, calcolati con un solver esatto sui dati CREA.", 11.5, True, GREEN),
    ("Porzioni a multipli di 5-10 g · frutti grandi a pezzi · pesi a crudo.", 11, False, GREY),
], sp=3)

# ---------- GIORNI ----------
week = PLAN["week"]
for day in ["Lunedi", "Martedi", "Mercoledi", "Giovedi", "Venerdi", "Sabato", "Domenica"]:
    r = week[day]; t = r["totals"]; bm = r["by_meal"]
    s = prs.slides.add_slide(BLANK); bg(s)
    # immagine a destra
    pic(s, DAY_IMG[day], 9.0, 0, 4.333, 7.5, f"day_{day}")
    # header
    text(s, 0.6, 0.35, 8.2, 0.7, [(day.upper(), 30, True, GREEN_D)])
    text(s, 0.62, 1.02, 8.2, 0.4, [(DAY_TAG[day], 15, False, ORANGE)])
    # badge macro
    bx, bw, gap = 0.6, 1.95, 0.12
    badge(s, bx + 0*(bw+gap), 1.55, bw, "KCAL", f"{t['kcal']:.0f}", GREEN)
    badge(s, bx + 1*(bw+gap), 1.55, bw, "PROTEINE", f"{t['protein_g']:.0f} g", GREEN_D)
    badge(s, bx + 2*(bw+gap), 1.55, bw, "CARBO", f"{t['carb_g']:.0f} g", ORANGE)
    badge(s, bx + 3*(bw+gap), 1.55, bw, "GRASSI", f"{t['fat_g']:.0f} g", RGBColor(0xB5,0x6A,0x2E))
    # card pasti 2x2
    cw, ch = 4.0, 1.95; cx0, cy0 = 0.6, 2.75; cgx, cgy = 0.2, 0.18
    meals = [m for m in MEAL_ORDER if m in bm]
    for i, m in enumerate(meals):
        col, row = i % 2, i // 2
        l = cx0 + col * (cw + cgx); tp = cy0 + row * (ch + cgy)
        box(s, l, tp, cw, ch, CARD, rounded=True)
        box(s, l, tp, 0.12, ch, GREEN if m != "cena" else ORANGE)  # barretta colore
        items = ", ".join(label(it) for it in bm[m]["items"])
        text(s, l + 0.22, tp + 0.12, cw - 0.35, 0.35,
             [(MEAL_LABEL[m] + f"   {bm[m]['kcal']:.0f} kcal", 12, True, GREEN_D)])
        text(s, l + 0.22, tp + 0.55, cw - 0.35, ch - 0.6, [(items, 11.5, False, INK)])
    if day == "Sabato":
        # nota sgarro al posto della cena (manca nei pasti controllati)
        l = cx0 + 1 * (cw + cgx); tp = cy0 + 1 * (ch + cgy)
        box(s, l, tp, cw, ch, RGBColor(0xF7,0xE3,0xC8), rounded=True)
        box(s, l, tp, 0.12, ch, ORANGE)
        text(s, l + 0.22, tp + 0.12, cw - 0.35, 0.35, [("CENA · SGARRO", 12, True, ORANGE)])
        text(s, l + 0.22, tp + 0.55, cw - 0.35, ch - 0.6,
             [("Pasto libero (~700-800 kcal): qui ci sta la pasta, la pizza o quello che vuoi.", 11.5, False, INK)])
    alt = r.get("cena_alt")
    if alt and alt.get("by_meal"):
        aitems = ", ".join(label(it) for it in alt["by_meal"]["items"])
        text(s, 0.6, 6.95, 8.2, 0.5,
             [(f"In alternativa alla cena — {alt['label']}: {aitems}", 10, False, GREEN_D)])

# ---------- LISTA SPESA ----------
s = prs.slides.add_slide(BLANK); bg(s)
pic(s, "veg", 9.6, 0, 3.733, 7.5, "shop")
text(s, 0.6, 0.4, 8.5, 0.7, [("Lista della spesa", 30, True, GREEN_D)])
text(s, 0.62, 1.05, 8.5, 0.4, [("Totale settimana · pesi a crudo", 14, False, ORANGE)])
shop = PLAN["shopping"]; bycat = {}
for n, g in shop.items():
    bycat.setdefault(DB.get(n, {}).get("category", "Altro"), []).append((n, g))
order = ["Carni fresche", "Prodotti della pesca", "Formaggi e latticini", "Latte e yogurt", "Uova",
         "Cereali e derivati", "Legumi", "Verdure e ortaggi", "Frutta", "Oli e grassi",
         "Frutta secca a guscio e semi oleaginosi", "Dolci"]
cats = [c for c in order if c in bycat]
half = (len(cats) + 1) // 2
cols = [cats[:half], cats[half:]]
for ci, colcats in enumerate(cols):
    runs = []
    for c in colcats:
        runs.append((c.upper(), 12, True, GREEN))
        s_items = ", ".join(
            (f"{short(n)} {g/1000:.1f} kg" if g >= 1000 else f"{short(n)} {int(g)} g")
            for n, g in sorted(bycat[c], key=lambda x: -x[1]))
        runs.append((s_items, 10.5, False, INK))
    text(s, 0.6 + ci * 4.45, 1.7, 4.3, 5.6, runs, sp=5)

# ---------- NOTE ----------
s = prs.slides.add_slide(BLANK); bg(s)
pic(s, "fruit", 9.6, 0, 3.733, 7.5, "notes")
text(s, 0.6, 0.4, 8.5, 0.7, [("Come usarla", 30, True, GREEN_D)])
notes = [
 ("Pesi a crudo", "Pasta, farro, couscous e carne/pesce si pesano CRUDI. Ceci e lenticchie sono in scatola, SCOLATI (gia' cotti). Tonno sgocciolato; salumi, formaggi e feta come stanno."),
 ("Frutti grandi a pezzi", "Banana e pesca sono indicate a numero (1, 2...), non a grammi."),
 ("Verdura libera", "Dove la verdura e' poca, aumentala quanto vuoi: incide pochissimo sulle calorie."),
 ("Tzatziki (mar/mer)", "Cetrioli + yogurt greco 0% + aglio, succo di limone, menta e un filo d'olio."),
 ("Sodio", "Tonno, feta, salumi e mozzarella sono salati: non aggiungere sale, sciacqua tonno e legumi in scatola."),
 ("Taratura", "Pesati 2 settimane allo stesso orario. Peso stabile = sei a target; se sale togli ~150 kcal, se scende aggiungi."),
]
runs = []
for tit, body in notes:
    runs.append((tit, 14, True, ORANGE))
    runs.append((body, 12, False, INK))
text(s, 0.6, 1.7, 8.6, 5.6, runs, sp=9)

prs.save(OUT)
print("SCRITTO", OUT, "-", len(prs.slides.__iter__.__self__._sldIdLst), "slide")
